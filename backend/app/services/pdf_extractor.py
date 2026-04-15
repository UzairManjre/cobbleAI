import pdfplumber
import docx
import pptx
from io import BytesIO

def extract_pdf(file_bytes: bytes):
    text = ""
    with pdfplumber.open(BytesIO(file_bytes)) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                text += t + "\n"
    return text

def extract_docx(file_bytes: bytes):
    doc = docx.Document(BytesIO(file_bytes))
    return "\n".join([p.text for p in doc.paragraphs])

def extract_pptx(file_bytes: bytes):
    presentation = pptx.Presentation(BytesIO(file_bytes))
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text(file_type: str, file_bytes: bytes) -> str:
    if file_type == "pdf":
        return extract_pdf(file_bytes)
    elif file_type == "docx":
        return extract_docx(file_bytes)
    elif file_type == "pptx":
        return extract_pptx(file_bytes)
    else:
        raise ValueError("Unsupported file type")
